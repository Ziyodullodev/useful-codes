from pptx import Presentation 
 
# Create a presentation object 
presentation = Presentation() 
 
# Title Slide 
slide = presentation.slides.add_slide(presentation.slide_layouts[0]) 
title = slide.shapes.title 
subtitle = slide.placeholders[1] 
title.text = "Python Dasturlashga Kirish" 
subtitle.text = "1-dars: Dasturlash asoslari va Python haqida tushunchalar" 
 
# Slide 1: Roadmap 
slide = presentation.slides.add_slide(presentation.slide_layouts[1]) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Dars Rejasi (Roadmap)" 
content.text = """\ 
1. Tanishuv 
2. Dasturlash haqida tushunchalar 
3. Kompyuterlar sozligi va hammada borligi 
4. Python haqida tushuncha 
5. Backend va frontend 
6. Python sohalari va imkoniyatlari 
7. Pythonni o'rnatish va onlayn platformalar""" 
 
# Slide 2: Dasturlash haqida 
slide = presentation.slides.add_slide(presentation.slide_layouts[1]) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Dasturlash haqida tushunchalar" 
content.text = """\ 
- Dasturlash - kompyuterga vazifalarni avtomatik bajarishni o'rgatish san'ati. 
- Kompyuter dasturchilarning buyruqlariga asoslangan holda ishlaydi. 
- Kundalik hayotimizda texnologiyaning asosiy qismi.""" 
 
# Slide 3: Python haqida tushuncha 
slide = presentation.slides.add_slide(presentation.slide_layouts[1]) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Python haqida tushuncha" 
content.text = """\ 
- Python - oddiy va qulay dasturlash tili. 
- 1991 yilda Guido van Rossum tomonidan yaratilgan. 
- Soddaligi sababli yangi boshlovchilar uchun ideal. 
- Juda ko'p sohalarda qo'llaniladi.""" 
 
# Slide 4: Backend va frontend 
slide = presentation.slides.add_slide(presentation.slide_layouts[1]) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Backend va Frontend" 
content.text = """\ 
- Frontend: Foydalanuvchi ko'radigan qismi (vazifa: sayt dizayni va interaktivligi). 
- Backend: Dastur ishlash tizimi (vazifa: ma'lumotlarni qayta ishlash va saqlash). 
- Backend foydalari: Kuchli logika, server boshqaruvi. 
- Backend zaifliklari: Ko'proq texnik bilim talab qiladi.""" 
 
# Slide 5: Python sohalari 
slide = presentation.slides.add_slide(presentation.slide_layouts[1]) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Python sohalari" 
content.text = """\ 
- Web dasturlash (backend) 
- Data Science 
- Telegram botlar 
- Sun'iy intellekt 
- Android va Desktop dasturlar""" 
 
# Slide 6: Hozir dasturchilarga talablar 
slide = presentation.slides.add_slide(presentation.slide_layouts[1]) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Hozir dasturchilarga talablar" 
content.text = """\ 
- Algoritmik fikrlash qobiliyati 
- Muammolarni yechish ko'nikmalari 
- Yangi texnologiyalarni o'rganishga qiziqish""" 
 
# Slide 7: Pythonni o'rnatish va onlayn ishlash 
slide = presentation.slides.add_slide(presentation.slide_layouts[1]) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Pythonni o'rnatish va onlayn ishlash" 
content.text = """\ 
1. Pythonni kompyuterga o'rnatish: python.org saytidan yuklab olish. 
2. Onlayn platformalar: 
   - Replit (replit.com) 
   - Google Colab (colab.research.google.com)""" 
 
# Save the presentation 
file_path = "/mnt/data/Python_Dars_1.pptx" 
presentation.save(file_path) 
 
file_path
